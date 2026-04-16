"""
Generate docs/presentation.pptx  –  XJCO3011 Coursework 1
Slides for 5-minute oral presentation + Q&A support.

Run from project root:  python docs/generate_pptx.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from lxml import etree

OUT = os.path.join(os.path.dirname(__file__), "presentation.pptx")

# ── Colour palette (Leeds-inspired deep blue + accents) ──────────────────────
BLUE_DARK   = RGBColor(0x0E, 0x2A, 0x5C)   # deep navy
BLUE_MID    = RGBColor(0x1A, 0x4A, 0x9E)   # university blue
BLUE_LIGHT  = RGBColor(0xD6, 0xE4, 0xF7)   # pale blue
GOLD        = RGBColor(0xF5, 0xA5, 0x00)   # gold accent
GREEN       = RGBColor(0x1A, 0x7A, 0x40)   # success green
RED         = RGBColor(0xC0, 0x39, 0x2B)   # error red
ORANGE      = RGBColor(0xE6, 0x7E, 0x22)   # warning orange
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT  = RGBColor(0xF4, 0xF6, 0xFA)
GRAY_MID    = RGBColor(0x7F, 0x8C, 0x8D)
BLACK       = RGBColor(0x1A, 0x1A, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helper utilities ─────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]   # completely blank


def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE=1
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = Pt(line_width or 1)
    else:
        line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, italic=False,
             color=BLACK, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name      = font_name
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font_size=16, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, space_before=0, font_name="Calibri",
             indent=None):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    if indent is not None:
        p.level = indent
    run = p.add_run()
    run.text = text
    run.font.name      = font_name
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = color
    return p


def header_bar(slide, title_text, subtitle_text=None):
    """Blue header bar across the top."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), fill_rgb=BLUE_DARK)
    add_text(slide, title_text,
             Inches(0.35), Inches(0.15), Inches(11), Inches(0.65),
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_text(slide, subtitle_text,
                 Inches(0.35), Inches(0.75), Inches(11), Inches(0.35),
                 font_size=14, bold=False, color=BLUE_LIGHT, align=PP_ALIGN.LEFT)
    # slide number placeholder at top-right
    return slide


def footer_bar(slide, left_text="XJCO3011 · Nutrition & Recipe Analytics API",
               right_text="University of Leeds / SWJTU"):
    add_rect(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), fill_rgb=BLUE_DARK)
    add_text(slide, left_text,
             Inches(0.2), Inches(7.12), Inches(7), Inches(0.3),
             font_size=9, color=BLUE_LIGHT)
    add_text(slide, right_text,
             Inches(10), Inches(7.12), Inches(3), Inches(0.3),
             font_size=9, color=BLUE_LIGHT, align=PP_ALIGN.RIGHT)


def bullet_box(slide, items, left, top, width, height,
               font_size=16, color=BLACK, bullet_char="▸", gap=4):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(gap)
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.size      = Pt(font_size)
        run.font.color.rgb = color
        run.font.name      = "Calibri"
    return txb


def colored_badge(slide, text, left, top, width, height,
                  bg_rgb=BLUE_MID, text_rgb=WHITE, font_size=13, bold=True):
    add_rect(slide, left, top, width, height, fill_rgb=bg_rgb)
    add_text(slide, text, left, top, width, height,
             font_size=font_size, bold=bold, color=text_rgb,
             align=PP_ALIGN.CENTER)


def section_divider(slide, number, title, color=BLUE_MID):
    """Full-slide section divider."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=BLUE_DARK)
    add_text(slide, number,
             Inches(4.5), Inches(2.5), Inches(4), Inches(1.2),
             font_size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, title,
             Inches(2), Inches(3.6), Inches(9.33), Inches(1),
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title slide."""
    slide = prs.slides.add_slide(blank_layout(prs))
    # Full background
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=BLUE_DARK)
    # Gold accent bar on left
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_rgb=GOLD)
    # Main title
    add_text(slide, "Nutrition & Recipe Analytics API",
             Inches(0.5), Inches(1.4), Inches(12), Inches(1.5),
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Subtitle
    add_text(slide, "XJCO3011 Web Services and Web Data  ·  Coursework 1",
             Inches(0.5), Inches(2.9), Inches(10), Inches(0.5),
             font_size=20, color=GOLD, align=PP_ALIGN.LEFT)
    # Description
    add_text(slide,
             "A fully functional RESTful API backed by 3 open food datasets:\n"
             "USDA FoodData Central · Food.com Recipes · Open Food Facts",
             Inches(0.5), Inches(3.5), Inches(10), Inches(1.0),
             font_size=16, color=BLUE_LIGHT, align=PP_ALIGN.LEFT)
    # Key stats row
    stats = [
        ("30", "HTTP Endpoints"),
        ("56",  "Passing Tests"),
        ("7,793", "Ingredients"),
        ("231,636", "Recipes"),
        ("10", "MCP Tools"),
    ]
    x = Inches(0.5)
    for val, label in stats:
        add_rect(slide, x, Inches(4.8), Inches(2.2), Inches(0.9), fill_rgb=BLUE_MID)
        add_text(slide, val, x, Inches(4.82), Inches(2.2), Inches(0.5),
                 font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, Inches(5.2), Inches(2.2), Inches(0.4),
                 font_size=11, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
        x += Inches(2.4)
    # Bottom info
    add_text(slide, "GitHub: github.com/Qhf7/XJCO3011-CW1",
             Inches(0.5), Inches(6.1), Inches(6), Inches(0.4),
             font_size=13, color=BLUE_LIGHT)
    add_text(slide, "Live API: web-production-e0934.up.railway.app",
             Inches(0.5), Inches(6.5), Inches(6), Inches(0.4),
             font_size=13, color=BLUE_LIGHT)
    add_text(slide, "Deadline: 21 April 2026",
             Inches(9.5), Inches(6.5), Inches(3.5), Inches(0.4),
             font_size=13, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_02_overview(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Project Overview", "What · Why · How")
    footer_bar(slide)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(5.9), fill_rgb=GRAY_LIGHT)

    # Left column: What is it?
    add_rect(slide, Inches(0.3), Inches(1.4), Inches(5.8), Inches(5.4), fill_rgb=WHITE)
    add_text(slide, "What is it?", Inches(0.5), Inches(1.5), Inches(5.4), Inches(0.4),
             font_size=16, bold=True, color=BLUE_DARK)
    add_rect(slide, Inches(0.3), Inches(1.88), Inches(5.8), Inches(0.05), fill_rgb=GOLD)

    items_left = [
        "RESTful API for nutrition analysis and recipe management",
        "Backed by 3 authoritative open datasets (USDA, Food.com, OFF)",
        "30 HTTP endpoints across 5 domains",
        "JWT-protected write operations (OAuth2 Password Flow)",
        "MCP server exposes 10 tools for AI assistants (Claude, GPT-4)",
        "Deployed live on Railway (cloud platform)",
    ]
    bullet_box(slide, items_left, Inches(0.5), Inches(2.0), Inches(5.4), Inches(4.6),
               font_size=15, color=BLACK)

    # Right column: Why this topic?
    add_rect(slide, Inches(6.9), Inches(1.4), Inches(6.1), Inches(5.4), fill_rgb=WHITE)
    add_text(slide, "Why this topic?", Inches(7.1), Inches(1.5), Inches(5.7), Inches(0.4),
             font_size=16, bold=True, color=BLUE_DARK)
    add_rect(slide, Inches(6.9), Inches(1.88), Inches(6.1), Inches(0.05), fill_rgb=GOLD)

    items_right = [
        "Nutrition data is universally relevant — everyone eats",
        "Rich, public-domain datasets available (USDA, Kaggle)",
        "Complex analytics justify advanced implementation",
        "Opportunity to invent novel metrics (NDS, difficulty score)",
        "MCP integration showcases cutting-edge AI interoperability",
        "Real-world utility: dietary tracking, meal planning, allergen safety",
    ]
    bullet_box(slide, items_right, Inches(7.1), Inches(2.0), Inches(5.7), Inches(4.6),
               font_size=15, color=BLACK)


def slide_03_stack(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Technology Stack", "Justified choices — not defaults")
    footer_bar(slide)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(5.9), fill_rgb=GRAY_LIGHT)

    cards = [
        (BLUE_MID,  "FastAPI",        "Python 3.12",   "Auto-generated OpenAPI docs, Pydantic v2 validation,\nasync support. Chosen over Django REST: zero-config\nSwagger UI is a coursework requirement."),
        (GREEN,     "SQLAlchemy 2.0", "ORM + SQLite",   "Type-annotated Mapped[T] models. SQLite: no server\nprocess, sub-ms queries on 231K recipes, WAL mode.\nMigrate to PostgreSQL with one env-var change."),
        (ORANGE,    "JWT Auth",       "python-jose",    "OAuth2 Password Flow, RFC 6750 Bearer tokens.\npasslib/bcrypt password hashing. Reusable FastAPI\ndependency injected into all protected routes."),
        (RGBColor(0x6C, 0x35, 0x9E), "FastMCP 3.0",    "MCP Server",     "Exposes 10 nutrition tools to AI assistants.\nsearch_recipes, get_allergens, compare_ingredients…\nTargets 70–79 'MCP-compatible advanced features' band."),
        (RGBColor(0xC0, 0x39, 0x2B), "Pytest",         "56 Tests",       "4 modules, in-memory SQLite (StaticPool).\nCovers auth, CRUD, analytics, service unit tests.\nAll 56 pass in < 4 seconds."),
    ]

    x_positions = [Inches(0.3), Inches(2.85), Inches(5.4), Inches(7.95), Inches(10.5)]
    for i, (color, name, tech, desc) in enumerate(cards):
        x = x_positions[i]
        add_rect(slide, x, Inches(1.35), Inches(2.4), Inches(5.55), fill_rgb=WHITE)
        add_rect(slide, x, Inches(1.35), Inches(2.4), Inches(0.55), fill_rgb=color)
        add_text(slide, name, x, Inches(1.38), Inches(2.4), Inches(0.35),
                 font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, tech, x, Inches(1.75), Inches(2.4), Inches(0.28),
                 font_size=11, color=GRAY_MID, align=PP_ALIGN.CENTER)
        add_text(slide, desc, Inches(x.inches + 0.1), Inches(2.1), Inches(2.2), Inches(4.6),
                 font_size=12, color=BLACK, wrap=True)


def slide_04_architecture(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Architecture", "Clean 4-layer design for testability and separation of concerns")
    footer_bar(slide)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(5.9), fill_rgb=GRAY_LIGHT)

    # Architecture diagram as blocks
    layers = [
        (GOLD,      "HTTP Layer",      "FastAPI Routers",   "app/routers/*.py",  "Thin controllers — receive requests, call services, return responses"),
        (BLUE_MID,  "Schemas",         "Pydantic v2",       "app/schemas/*.py",  "Request validation & response serialization. 422 on malformed input"),
        (GREEN,     "Services",        "Business Logic",    "app/services/*.py", "Allergen detection, NDS scoring, difficulty estimation — pure Python"),
        (ORANGE,    "Data Layer",      "SQLAlchemy 2.0",    "app/models/*.py",   "ORM table definitions. nutrition.db (456 MB) / nutrition_demo.db (11.6 MB)"),
    ]

    y = Inches(1.45)
    for bg, title, subtitle, path, desc in layers:
        add_rect(slide, Inches(0.3), y, Inches(12.6), Inches(1.1), fill_rgb=WHITE)
        add_rect(slide, Inches(0.3), y, Inches(0.15), Inches(1.1), fill_rgb=bg)
        add_text(slide, title, Inches(0.6), y + Pt(6), Inches(2.2), Inches(0.5),
                 font_size=14, bold=True, color=BLUE_DARK)
        add_text(slide, subtitle, Inches(0.6), y + Pt(30), Inches(2.2), Inches(0.35),
                 font_size=11, color=GRAY_MID)
        # path badge
        add_rect(slide, Inches(2.9), y + Pt(10), Inches(2.3), Inches(0.38), fill_rgb=GRAY_LIGHT)
        add_text(slide, path, Inches(2.9), y + Pt(10), Inches(2.3), Inches(0.38),
                 font_size=11, color=BLUE_MID, align=PP_ALIGN.CENTER, font_name="Courier New")
        add_text(slide, desc, Inches(5.4), y + Pt(10), Inches(7.5), Inches(0.55),
                 font_size=13, color=BLACK)
        y += Inches(1.18)

    # MCP Server side note
    add_rect(slide, Inches(0.3), Inches(5.9), Inches(12.6), Inches(0.85), fill_rgb=RGBColor(0xEE, 0xE5, 0xF9))
    add_rect(slide, Inches(0.3), Inches(5.9), Inches(0.15), Inches(0.85), fill_rgb=RGBColor(0x6C, 0x35, 0x9E))
    add_text(slide, "MCP Server  (app/mcp_server.py)",
             Inches(0.6), Inches(5.93), Inches(3.5), Inches(0.4),
             font_size=13, bold=True, color=RGBColor(0x6C, 0x35, 0x9E))
    add_text(slide, "Runs alongside FastAPI · exposes 10 tools to Claude/GPT-4 via Model Context Protocol",
             Inches(4.0), Inches(5.93), Inches(8.9), Inches(0.4),
             font_size=13, color=BLACK)


def slide_05_datasets(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Data Sources", "Three authoritative open datasets")
    footer_bar(slide)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(5.9), fill_rgb=GRAY_LIGHT)

    datasets = [
        (BLUE_MID, "USDA FoodData Central", "SR Legacy 2018 · Public Domain",
         ["7,793 foods with precise per-100g measurements",
          "474 nutrient types: energy, macros, vitamins, minerals",
          "Source of truth for ingredient nutrition",
          "URL: fdc.nal.usda.gov"],
         "7,793\nfoods"),
        (GREEN, "Food.com Recipes", "Kaggle · CC BY-SA 3.0",
         ["231,636 real-world recipes",
          "Ingredients list, steps, cooking time, tags",
          "Per-serving %DV nutritional estimates",
          "Used for recipe CRUD + analytics"],
         "231,636\nrecipes"),
        (ORANGE, "Open Food Facts", "Cache-first · ODbL",
         ["Real-time lookup by barcode or product name",
          "Allergen data, Nutri-Score grade (A–E)",
          "Results cached locally (off_product table)",
          "Keeps deployment DB at 11.6 MB"],
         "Live\nAPI"),
    ]

    x_positions = [Inches(0.3), Inches(4.55), Inches(8.8)]
    for i, (color, name, sub, points, stat) in enumerate(datasets):
        x = x_positions[i]
        w = Inches(4.0)
        add_rect(slide, x, Inches(1.45), w, Inches(5.4), fill_rgb=WHITE)
        add_rect(slide, x, Inches(1.45), w, Inches(0.8), fill_rgb=color)
        add_text(slide, stat, x + Inches(3.0), Inches(1.48), Inches(0.9), Inches(0.72),
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, name, x + Inches(0.15), Inches(1.5), Inches(2.8), Inches(0.4),
                 font_size=15, bold=True, color=WHITE)
        add_text(slide, sub, x + Inches(0.15), Inches(1.85), w - Inches(0.3), Inches(0.3),
                 font_size=11, color=GRAY_MID)
        bullet_box(slide, points,
                   x + Inches(0.15), Inches(2.3), w - Inches(0.3), Inches(4.3),
                   font_size=13, color=BLACK)


def slide_06_innovations(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Innovation Features", "Three proprietary algorithms beyond basic CRUD")
    footer_bar(slide)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(5.9), fill_rgb=GRAY_LIGHT)

    innovations = [
        (BLUE_MID, "🎯 Nutrient Density Score",
         "NDS = mean(protein, fibre, Vit-C, Ca, Fe scores)",
         ["Measures nutrition value per 100 kcal", "Graded A–E (like Nutri-Score)", "Spinach: Grade A  ·  Sugar: Grade E", "Endpoint: GET /ingredients/{id}/nutrient-density"],
         "A–E\nGrade"),
        (GREEN, "⭐ Difficulty Estimator",
         "Multi-factor 1–5 star scoring",
         ["Factors: steps + ingredients + time", "Technique detection: 30 advanced keywords", "(sous vide, deglaze, emulsify…)", "Endpoint: GET /recipes/{id}/difficulty"],
         "1–5\nStars"),
        (ORANGE, "⚠️ FDA Allergen Engine",
         "9 major FDA allergens + sulfites",
         ["Keyword matching across ingredient list", "Returns per-allergen ingredient breakdown", "safe_for labels: peanut-free, egg-free…", "Endpoint: GET /recipes/{id}/allergens"],
         "9+1\nAllergens"),
    ]

    x_positions = [Inches(0.3), Inches(4.55), Inches(8.8)]
    for i, (color, title, formula, points, badge) in enumerate(innovations):
        x = x_positions[i]
        w = Inches(4.0)
        add_rect(slide, x, Inches(1.45), w, Inches(5.4), fill_rgb=WHITE)
        add_rect(slide, x, Inches(1.45), w, Inches(0.75), fill_rgb=color)
        add_text(slide, badge, x + Inches(2.95), Inches(1.48), Inches(0.95), Inches(0.66),
                 font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(0.15), Inches(1.5), Inches(2.75), Inches(0.42),
                 font_size=14, bold=True, color=WHITE)
        add_text(slide, formula,
                 x + Inches(0.15), Inches(2.22), w - Inches(0.3), Inches(0.38),
                 font_size=12, italic=True, color=BLUE_MID, font_name="Courier New")
        bullet_box(slide, points,
                   x + Inches(0.15), Inches(2.65), w - Inches(0.3), Inches(4.0),
                   font_size=13, color=BLACK)


def slide_07_endpoints(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "API Endpoints", "30 endpoints across 5 domains")
    footer_bar(slide)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(5.9), fill_rgb=GRAY_LIGHT)

    domains = [
        (BLUE_MID, "🔐 Auth",     "2",  ["POST /auth/register", "POST /auth/login"]),
        (GREEN,    "🥦 Ingredients","8", ["GET /ingredients", "POST /ingredients", "GET /ingredients/{id}", "PUT /ingredients/{id}", "DELETE /ingredients/{id}", "GET …/nutrition", "GET …/nutrient-density", "GET /ingredients/categories"]),
        (ORANGE,   "🍽️ Recipes",  "9",  ["GET /recipes", "POST /recipes", "GET /recipes/{id}", "PUT /recipes/{id}", "DELETE /recipes/{id}", "GET …/allergens", "GET …/difficulty", "GET …/nutrition", "GET /recipes/by-ingredients"]),
        (RGBColor(0x6C, 0x35, 0x9E), "📊 Analytics", "7", ["POST /analytics/nutrition/calculate", "POST /analytics/ingredients/compare", "POST /analytics/meal-plan/analyze", "GET /analytics/calories/distribution", "GET /analytics/top-ingredients", "GET /analytics/macro-trends", "GET /analytics/allergen-stats"]),
        (RGBColor(0x16, 0x7A, 0x8A), "🌍 Open Food Facts", "2+2", ["GET /food/barcode/{barcode}", "GET /food/search", "GET / (health)", "GET /stats"]),
    ]

    y_positions = [Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4)]
    col_widths = [Inches(2.0), Inches(3.5), Inches(3.0), Inches(3.5), Inches(1.0)]
    x_positions = [Inches(0.2), Inches(2.3), Inches(5.9), Inches(9.0), Inches(11.9)]
    # Actually do a simple layout - 5 columns
    x = Inches(0.2)
    col_w_list = [Inches(2.5), Inches(3.2), Inches(3.2), Inches(3.5), Inches(0.6)]

    for i, (color, domain, count, endpoints) in enumerate(domains):
        w = col_w_list[i]
        add_rect(slide, x, Inches(1.38), w, Inches(5.55), fill_rgb=WHITE)
        add_rect(slide, x, Inches(1.38), w, Inches(0.65), fill_rgb=color)
        add_text(slide, domain, x + Inches(0.1), Inches(1.42), w - Inches(0.2), Inches(0.35),
                 font_size=13, bold=True, color=WHITE)
        # count badge
        add_rect(slide, x + w - Inches(0.5), Inches(1.5), Inches(0.4), Inches(0.35),
                 fill_rgb=GOLD)
        add_text(slide, count, x + w - Inches(0.5), Inches(1.5), Inches(0.4), Inches(0.35),
                 font_size=12, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        # endpoints list
        for j, ep in enumerate(endpoints):
            yy = Inches(2.15) + j * Inches(0.48)
            add_text(slide, ep, x + Inches(0.1), yy, w - Inches(0.2), Inches(0.42),
                     font_size=10.5, color=BLACK, font_name="Courier New")
        x += w + Inches(0.1)


def slide_08_mcp(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "MCP Integration", "AI-native interface — Claude & GPT-4 can call your API directly")
    footer_bar(slide)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(5.9), fill_rgb=GRAY_LIGHT)

    # Left: explanation
    add_rect(slide, Inches(0.3), Inches(1.4), Inches(5.8), Inches(5.5), fill_rgb=WHITE)
    add_text(slide, "What is MCP?", Inches(0.5), Inches(1.5), Inches(5.4), Inches(0.4),
             font_size=16, bold=True, color=BLUE_DARK)
    add_rect(slide, Inches(0.3), Inches(1.9), Inches(5.8), Inches(0.04), fill_rgb=GOLD)

    bullet_box(slide, [
        "Model Context Protocol — open standard by Anthropic (2024)",
        "Allows AI assistants to call external tools as functions",
        "FastMCP 3.0 wraps FastAPI functions as MCP tools",
        "AI can query nutrition data in natural language",
        "Targets 70–79 marking band: 'MCP-compatible advanced features'",
    ], Inches(0.5), Inches(2.0), Inches(5.4), Inches(4.5), font_size=14, color=BLACK)

    # Right: 10 tools
    add_rect(slide, Inches(6.9), Inches(1.4), Inches(6.1), Inches(5.5), fill_rgb=WHITE)
    add_text(slide, "10 MCP Tools  (app/mcp_server.py)",
             Inches(7.1), Inches(1.5), Inches(5.7), Inches(0.4),
             font_size=16, bold=True, color=BLUE_DARK)
    add_rect(slide, Inches(6.9), Inches(1.9), Inches(6.1), Inches(0.04), fill_rgb=GOLD)

    tools = [
        ("search_recipes",         "Find recipes by name/ingredients"),
        ("get_recipe_allergens",   "Detect allergens in a recipe"),
        ("get_recipe_difficulty",  "Get 1–5 difficulty score"),
        ("search_ingredients",     "Search USDA ingredient database"),
        ("get_ingredient_nutrition","Per-100g macro & micro nutrients"),
        ("get_nutrient_density",   "NDS score A–E grade"),
        ("compare_ingredients",    "Side-by-side nutrition comparison"),
        ("calculate_nutrition",    "Custom meal nutrition total"),
        ("analyze_meal_plan",      "Full-day DRI coverage analysis"),
        ("get_api_stats",          "Live database statistics"),
    ]

    y = Inches(2.05)
    for tool, desc in tools:
        add_rect(slide, Inches(7.1), y, Inches(5.6), Inches(0.44), fill_rgb=GRAY_LIGHT)
        add_text(slide, tool, Inches(7.2), y + Pt(3), Inches(2.6), Inches(0.35),
                 font_size=11, bold=True, color=BLUE_MID, font_name="Courier New")
        add_text(slide, desc, Inches(9.9), y + Pt(3), Inches(2.7), Inches(0.35),
                 font_size=11, color=BLACK)
        y += Inches(0.48)


def slide_09_version_control(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Version Control", "15 commits · consistent incremental development history")
    footer_bar(slide)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(5.9), fill_rgb=GRAY_LIGHT)

    commits = [
        ("#1", "Initial project structure and FastAPI setup",         "Core skeleton, dependencies, app/main.py"),
        ("#2", "Add dataset documentation and data import scripts",   "USDA, Food.com, OFF import scripts + Data/README.md"),
        ("#3", "Add database models and schemas",                     "SQLAlchemy ORM models, Pydantic v2 schemas"),
        ("#4", "Implement authentication module",                     "JWT auth, OAuth2 Password Flow, passlib/bcrypt"),
        ("#5", "Add ingredients and recipes CRUD endpoints",          "app/routers/ingredients.py + recipes.py"),
        ("#6", "Implement analytics services and endpoints",          "NDS, difficulty estimator, allergen engine"),
        ("#7", "Integrate Open Food Facts API",                       "Cache-first barcode/name lookup"),
        ("#8", "Add MCP server with 10 tools",                        "FastMCP 3.0, app/mcp_server.py"),
        ("#9", "Add comprehensive test suite",                        "56 pytest tests, StaticPool in-memory DB"),
        ("#10","Add pre-built demo database",                         "nutrition_demo.db (11.6 MB), create_demo_db.py"),
        ("#11","Configure Railway deployment",                         "Procfile, railway.toml, render.yaml, wsgi.py"),
        ("#12","Add technical report and API documentation",           "docs/technical_report.pdf, api_documentation.pdf"),
        ("#13","Add GenAI appendix",                                   "docs/genai_appendix.pdf (161 messages, 7 stages)"),
        ("#14","Fix ingredient CRUD and Pydantic v2 compatibility",    "FOREIGN KEY bug, ConfigDict migration"),
        ("#15","Docs: update tech report & API docs with Railway URL", "Lessons Learned, References, live base URL"),
    ]

    # Two-column layout
    y_left  = Inches(1.4)
    y_right = Inches(1.4)
    for i, (num, title, detail) in enumerate(commits):
        col = i % 2
        x = Inches(0.3) if col == 0 else Inches(6.8)
        y = y_left       if col == 0 else y_right

        add_rect(slide, x, y, Inches(6.3), Inches(0.43), fill_rgb=WHITE)
        # Number badge
        badge_color = BLUE_MID if int(num[1:]) <= 7 else GREEN
        add_rect(slide, x, y, Inches(0.45), Inches(0.43), fill_rgb=badge_color)
        add_text(slide, num, x, y, Inches(0.45), Inches(0.43),
                 font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(0.5), y + Pt(2), Inches(4.5), Inches(0.24),
                 font_size=10.5, bold=True, color=BLACK)
        add_text(slide, detail, x + Inches(0.5), y + Pt(18), Inches(5.7), Inches(0.2),
                 font_size=9, color=GRAY_MID)

        if col == 0:
            y_left  += Inches(0.49)
        else:
            y_right += Inches(0.49)

    # GitHub link
    add_text(slide, "🔗  github.com/Qhf7/XJCO3011-CW1",
             Inches(0.3), Inches(7.0), Inches(8), Inches(0.35),
             font_size=13, bold=True, color=BLUE_MID)


def slide_10_testing(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Testing", "56 tests · 4 modules · 100% pass rate")
    footer_bar(slide)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(5.9), fill_rgb=GRAY_LIGHT)

    # Test modules table
    modules = [
        (BLUE_MID, "test_auth.py",        "8",  "Registration, login, JWT token validation, access control, invalid credentials"),
        (GREEN,    "test_ingredients.py", "12", "Full CRUD, nutrition lookup, NDS grading, search, USDA read-only protection"),
        (ORANGE,   "test_recipes.py",     "17", "Full CRUD, allergen detection, difficulty scoring, nutrition scaling, smart filters"),
        (RGBColor(0x6C, 0x35, 0x9E), "test_analytics.py", "19", "Service unit tests (allergen/difficulty/NDS), analytics endpoints, nutrition calculator, meal plan"),
    ]

    y = Inches(1.4)
    for color, mod, count, desc in modules:
        add_rect(slide, Inches(0.3), y, Inches(12.6), Inches(1.1), fill_rgb=WHITE)
        add_rect(slide, Inches(0.3), y, Inches(0.12), Inches(1.1), fill_rgb=color)
        add_text(slide, mod, Inches(0.55), y + Pt(5), Inches(2.8), Inches(0.4),
                 font_size=14, bold=True, color=BLACK, font_name="Courier New")
        add_text(slide, f"{count} tests", Inches(0.55), y + Pt(28), Inches(2.8), Inches(0.3),
                 font_size=12, color=color, bold=True)
        add_text(slide, desc, Inches(3.5), y + Pt(10), Inches(9.3), Inches(0.6),
                 font_size=13, color=BLACK)
        y += Inches(1.2)

    # Key technical note
    add_rect(slide, Inches(0.3), Inches(6.2), Inches(12.6), Inches(0.65),
             fill_rgb=RGBColor(0xE8, 0xF5, 0xE9))
    add_text(slide,
             "⚙️  Test isolation:  SQLAlchemy StaticPool ensures all test sessions share one in-memory connection  "
             "·  Production database never modified  ·  All 56 tests complete in < 4 seconds",
             Inches(0.5), Inches(6.25), Inches(12.2), Inches(0.5),
             font_size=12.5, color=GREEN, bold=True)


def slide_11_deployment(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Deployment", "Live on Railway — Professional cloud deployment")
    footer_bar(slide)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(5.9), fill_rgb=GRAY_LIGHT)

    # Left: deployment info
    add_rect(slide, Inches(0.3), Inches(1.4), Inches(6.2), Inches(5.5), fill_rgb=WHITE)
    add_text(slide, "Railway (railway.app)", Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.4),
             font_size=18, bold=True, color=BLUE_DARK)
    add_rect(slide, Inches(0.3), Inches(1.9), Inches(6.2), Inches(0.04), fill_rgb=GOLD)

    bullet_box(slide, [
        "Automatic deployment from GitHub main branch",
        "Start command: uvicorn app.main:app --host 0.0.0.0 --port $PORT",
        "Health check: GET / (returns 200 OK)",
        "2 vCPU · 1 GB RAM — handles full 231K recipe dataset",
        "nutrition_demo.db (11.6 MB) tracked in repository",
        "Full database at runtime: config.py auto-selects available file",
    ], Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.8),
       font_size=14, color=BLACK)

    # Right: URLs and config files
    add_rect(slide, Inches(6.8), Inches(1.4), Inches(6.2), Inches(5.5), fill_rgb=WHITE)
    add_text(slide, "Live URLs", Inches(7.0), Inches(1.5), Inches(6.0), Inches(0.4),
             font_size=18, bold=True, color=BLUE_DARK)
    add_rect(slide, Inches(6.8), Inches(1.9), Inches(6.2), Inches(0.04), fill_rgb=GOLD)

    urls = [
        ("API Root",        "web-production-e0934.up.railway.app"),
        ("Swagger UI",      "web-production-e0934.up.railway.app/docs"),
        ("ReDoc",           "web-production-e0934.up.railway.app/redoc"),
        ("OpenAPI JSON",    "web-production-e0934.up.railway.app/openapi.json"),
        ("Health Check",    "web-production-e0934.up.railway.app/ → 200 OK"),
    ]
    y = Inches(2.1)
    for label, url in urls:
        add_rect(slide, Inches(6.9), y, Inches(5.9), Inches(0.6), fill_rgb=GRAY_LIGHT)
        add_text(slide, label, Inches(7.0), y + Pt(2), Inches(1.9), Inches(0.28),
                 font_size=11, bold=True, color=GRAY_MID)
        add_text(slide, url, Inches(7.0), y + Pt(18), Inches(5.7), Inches(0.3),
                 font_size=11, color=BLUE_MID, font_name="Courier New")
        y += Inches(0.67)

    add_text(slide, "Config files: Procfile · railway.toml · render.yaml · wsgi.py",
             Inches(7.0), Inches(5.9), Inches(5.8), Inches(0.4),
             font_size=12, color=GRAY_MID, italic=True)


def slide_12_api_docs(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "API Documentation", "Auto-generated Swagger UI + 12-page PDF reference")
    footer_bar(slide)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(5.9), fill_rgb=GRAY_LIGHT)

    # Left: documentation features
    add_rect(slide, Inches(0.3), Inches(1.4), Inches(5.8), Inches(5.5), fill_rgb=WHITE)
    add_text(slide, "What's documented?", Inches(0.5), Inches(1.5), Inches(5.4), Inches(0.4),
             font_size=16, bold=True, color=BLUE_DARK)
    add_rect(slide, Inches(0.3), Inches(1.9), Inches(5.8), Inches(0.04), fill_rgb=GOLD)

    bullet_box(slide, [
        "All 30 endpoints with parameters and response formats",
        "Authentication flow (3-step: register → login → Bearer token)",
        "Full Error Code Reference (400, 401, 403, 404, 422, 500)",
        "10 JSON request/response examples (pages 8–12 of PDF)",
        "Swagger UI: interactive testing in browser",
        "ReDoc: clean reading-friendly format",
        "PDF exported from live OpenAPI 3.1.0 spec",
    ], Inches(0.5), Inches(2.0), Inches(5.4), Inches(4.8),
       font_size=14, color=BLACK)

    # Right: example JSON
    add_rect(slide, Inches(6.9), Inches(1.4), Inches(6.1), Inches(5.5), fill_rgb=WHITE)
    add_text(slide, "Example: GET /recipes/{id}/allergens",
             Inches(7.1), Inches(1.5), Inches(5.7), Inches(0.4),
             font_size=14, bold=True, color=BLUE_DARK)
    add_rect(slide, Inches(6.9), Inches(1.9), Inches(6.1), Inches(0.04), fill_rgb=GOLD)

    code = '''{
  "recipe_id": 42,
  "recipe_name": "Classic Mac and Cheese",
  "allergens_detected": [
    {
      "allergen": "dairy",
      "ingredients": ["cheddar", "butter", "milk"]
    },
    {
      "allergen": "gluten",
      "ingredients": ["elbow macaroni"]
    }
  ],
  "safe_for": ["peanut-free", "egg-free", "soy-free"],
  "allergen_free": false
}'''
    add_rect(slide, Inches(7.0), Inches(2.05), Inches(5.9), Inches(4.7),
             fill_rgb=RGBColor(0x1E, 0x1E, 0x2E))
    add_text(slide, code, Inches(7.1), Inches(2.1), Inches(5.7), Inches(4.6),
             font_size=12, color=RGBColor(0xA6, 0xE2, 0x2E),
             font_name="Courier New", wrap=True)


def slide_13_genai(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "Generative AI Usage", "GREEN light assessment — AI as primary thought partner")
    footer_bar(slide)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(5.9), fill_rgb=GRAY_LIGHT)

    stages = [
        (BLUE_MID,  "Dataset Discovery",    "AI analysed CSV schemas, identified 12 key USDA nutrient IDs from 474 available"),
        (BLUE_MID,  "Architecture Design",  "Explored FastAPI vs Django trade-offs; co-designed 4-layer architecture"),
        (GREEN,     "Algorithm Co-design",  "NDS formula and difficulty estimator: I specified intent, AI proposed maths, I validated"),
        (GREEN,     "Code Generation",      "AI generated initial modules; each reviewed, tested, refined (StaticPool, bcrypt fixes)"),
        (ORANGE,    "Debugging",            "Diagnosed starlette 1.0.0 conflict, SQLite VACUUM-in-transaction, pytest StaticPool"),
        (ORANGE,    "Documentation",        "AI drafted README, technical report, Swagger descriptions — all verified for accuracy"),
    ]

    y = Inches(1.4)
    for color, stage, desc in stages:
        add_rect(slide, Inches(0.3), y, Inches(12.6), Inches(0.78), fill_rgb=WHITE)
        add_rect(slide, Inches(0.3), y, Inches(0.12), Inches(0.78), fill_rgb=color)
        add_text(slide, stage, Inches(0.55), y + Pt(4), Inches(2.3), Inches(0.35),
                 font_size=13, bold=True, color=BLUE_DARK)
        add_text(slide, desc, Inches(2.9), y + Pt(10), Inches(10.0), Inches(0.45),
                 font_size=13, color=BLACK)
        y += Inches(0.85)

    # Bottom: grade band claim
    add_rect(slide, Inches(0.3), Inches(6.35), Inches(12.6), Inches(0.6),
             fill_rgb=RGBColor(0xFFF3, 0xCD00, 0x0000) if False else RGBColor(0xFF, 0xF3, 0xCD))
    add_text(slide,
             "🏆  Grade Band Target: 80–89 (Excellent)  —  "
             "\"High level use of GenAI to aid creative thinking and solution exploration\"",
             Inches(0.5), Inches(6.38), Inches(12.2), Inches(0.5),
             font_size=13, bold=True, color=RGBColor(0x85, 0x60, 0x04))


def slide_14_deliverables(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    header_bar(slide, "All Deliverables", "Everything submitted and accessible")
    footer_bar(slide)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(5.9), fill_rgb=GRAY_LIGHT)

    deliverables = [
        (BLUE_MID, "📁 Code Repository",
         "github.com/Qhf7/XJCO3011-CW1",
         ["Public GitHub repository", "15 meaningful commits", "README.md with setup instructions", "All code runnable locally and on Railway"]),
        (GREEN, "📄 API Documentation",
         "docs/api_documentation.pdf  (12 pages)",
         ["All 30 endpoints documented", "Auth flow + error codes", "10 JSON request/response examples", "Swagger UI: /docs   ReDoc: /redoc"]),
        (ORANGE, "📋 Technical Report",
         "docs/technical_report.pdf  (4 pages)",
         ["Stack justification + architecture", "Testing + Lessons Learned", "Limitations + Future Work", "GenAI declaration + References"]),
        (RGBColor(0x6C, 0x35, 0x9E), "📊 Presentation Slides",
         "docs/presentation.pptx  (this file)",
         ["Version control practices", "API documentation overview", "Technical report highlights", "All deliverables covered"]),
    ]

    x_positions = [Inches(0.3), Inches(3.6), Inches(6.9), Inches(10.2)]
    for i, (color, title, path, points) in enumerate(deliverables):
        x = x_positions[i]
        w = Inches(3.1)
        add_rect(slide, x, Inches(1.4), w, Inches(5.55), fill_rgb=WHITE)
        add_rect(slide, x, Inches(1.4), w, Inches(0.72), fill_rgb=color)
        add_text(slide, title, x + Inches(0.15), Inches(1.45), w - Inches(0.3), Inches(0.4),
                 font_size=14, bold=True, color=WHITE)
        add_text(slide, path, x + Inches(0.1), Inches(2.14), w - Inches(0.2), Inches(0.32),
                 font_size=10, color=BLUE_MID, font_name="Courier New")
        add_rect(slide, x, Inches(2.44), w, Inches(0.02), fill_rgb=BLUE_LIGHT)
        bullet_box(slide, points,
                   x + Inches(0.1), Inches(2.5), w - Inches(0.2), Inches(4.2),
                   font_size=12.5, color=BLACK, bullet_char="✓")

    # GenAI appendix note
    add_rect(slide, Inches(0.3), Inches(7.0), Inches(12.6), Inches(0.35), fill_rgb=BLUE_LIGHT)
    add_text(slide,
             "📎  GenAI Appendix:  docs/genai_appendix.pdf  —  "
             "161 AI messages across 7 development stages",
             Inches(0.5), Inches(7.02), Inches(12.2), Inches(0.3),
             font_size=12, color=BLUE_DARK, bold=True)


def slide_15_demo(prs):
    """Final slide: live demo prompt + Q&A."""
    slide = prs.slides.add_slide(blank_layout(prs))
    # Full dark background
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=BLUE_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_rgb=GOLD)

    add_text(slide, "Live Demo & Q&A",
             Inches(0.5), Inches(0.8), Inches(12), Inches(1.2),
             font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Demo steps
    steps = [
        ("1", "Open Swagger UI",      "web-production-e0934.up.railway.app/docs"),
        ("2", "Register & Login",     "POST /auth/register → POST /auth/login → copy token"),
        ("3", "Search recipes",       "GET /recipes?q=pasta → browse 231K real recipes"),
        ("4", "Check allergens",      "GET /recipes/{id}/allergens → FDA 9-allergen detection"),
        ("5", "NDS score",            "GET /ingredients/{id}/nutrient-density → Grade A–E"),
        ("6", "Nutrition calculator", "POST /analytics/nutrition/calculate → custom meal totals"),
    ]

    y = Inches(2.1)
    for num, title, desc in steps:
        add_rect(slide, Inches(0.5), y, Inches(0.45), Inches(0.5),
                 fill_rgb=GOLD)
        add_text(slide, num, Inches(0.5), y, Inches(0.45), Inches(0.5),
                 font_size=16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        add_text(slide, title, Inches(1.1), y + Pt(2), Inches(3.5), Inches(0.28),
                 font_size=14, bold=True, color=WHITE)
        add_text(slide, desc, Inches(1.1), y + Pt(22), Inches(11), Inches(0.24),
                 font_size=12, color=BLUE_LIGHT, font_name="Courier New")
        y += Inches(0.6)

    add_text(slide,
             "GitHub: github.com/Qhf7/XJCO3011-CW1     "
             "Live API: web-production-e0934.up.railway.app",
             Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
             font_size=13, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
def main():
    print(f"Generating {OUT} ...")
    prs = new_prs()

    slide_01_title(prs)
    slide_02_overview(prs)
    slide_03_stack(prs)
    slide_04_architecture(prs)
    slide_05_datasets(prs)
    slide_06_innovations(prs)
    slide_07_endpoints(prs)
    slide_08_mcp(prs)
    slide_09_version_control(prs)
    slide_10_testing(prs)
    slide_11_deployment(prs)
    slide_12_api_docs(prs)
    slide_13_genai(prs)
    slide_14_deliverables(prs)
    slide_15_demo(prs)

    prs.save(OUT)
    size = os.path.getsize(OUT) / 1024
    print(f"Done: {OUT}  ({size:.0f} KB,  {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
